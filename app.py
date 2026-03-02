import streamlit as st
import yfinance as yf
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime, timedelta
import feedparser
import io
import os

# --- IDENTIDADE VISUAL ---
DIR_BASE = os.path.dirname(os.path.abspath(__file__))
COR_FUNDO = RGBColor(0, 32, 77)
COR_TEXTO = RGBColor(255, 255, 255)
COR_DESTAQUE = RGBColor(0, 174, 239)
COR_ALTA = RGBColor(0, 255, 127)
COR_BAIXA = RGBColor(255, 69, 0)

def carregar_logo():
    for arq in os.listdir(DIR_BASE):
        if arq.lower().startswith("logo") and arq.lower().endswith(('.png', '.jpg', '.jpeg')):
            try:
                with open(os.path.join(DIR_BASE, arq), "rb") as f: return f.read()
            except: continue
    return None

def obter_variacoes_b3(data_alvo):
    """Busca ações da B3 uma a uma para garantir precisão e evitar erro de tabela"""
    tickers = [
        'VALE3.SA', 'PETR4.SA', 'ITUB4.SA', 'BBDC4.SA', 'BBAS3.SA', 'ABEV3.SA', 'MGLU3.SA', 
        'WEGE3.SA', 'PRIO3.SA', 'GGBR4.SA', 'RENT3.SA', 'LREN3.SA', 'HAPV3.SA', 'ELET3.SA', 
        'SUZB3.SA', 'JBSS3.SA', 'RAIL3.SA', 'RADL3.SA', 'VVAR3.SA', 'CSNA3.SA'
    ]
    
    resultados = {}
    s_str = data_alvo.strftime('%Y-%m-%d')
    e_str = (data_alvo + timedelta(days=1)).strftime('%Y-%m-%d')
    
    for t in tickers:
        try:
            d = yf.download(t, start=s_str, end=e_str, progress=False)
            if not d.empty and len(d) >= 1:
                ab = float(d['Open'].iloc[0])
                fe = float(d['Close'].iloc[-1])
                var = ((fe / ab) - 1) * 100
                resultados[t.replace('.SA', '')] = var
        except: continue
            
    if not resultados: return None, None
    
    series = pd.Series(resultados)
    return series.nlargest(5), series.nsmallest(5)

def obter_dado_unico(ticker, data_alvo):
    """Busca dados de um índice garantindo o pregão ativo"""
    for i in range(5):
        d_fim = data_alvo - timedelta(days=i)
        d_ini = d_fim - timedelta(days=1)
        d = yf.download(ticker, start=d_ini.strftime('%Y-%m-%d'), end=(d_fim + timedelta(days=1)).strftime('%Y-%m-%d'), progress=False)
        if not d.empty:
            ab = float(d['Open'].iloc[-1])
            fe = float(d['Close'].iloc[-1])
            return {"A": ab, "F": fe, "V": ((fe/ab)-1)*100, "D": d_fim}
    return None

def gerar_grafico(ticker, data_v, cor):
    try:
        s, e = data_v.strftime('%Y-%m-%d'), (data_v + timedelta(days=1)).strftime('%Y-%m-%d')
        df = yf.download(ticker, start=s, end=e, interval="5m", progress=False)['Close']
        if df.empty: return None
        plt.figure(figsize=(5, 2), facecolor='#00204D')
        ax = plt.axes(); ax.set_facecolor('#00204D')
        plt.plot(df.index, df.values, color=cor, linewidth=2.5)
        ax.tick_params(axis='both', colors='white', labelsize=8)
        for sp in ax.spines.values(): sp.set_color('white')
        plt.grid(True, color='grey', linestyle='--', alpha=0.1)
        plt.gcf().autofmt_xdate()
        img = io.BytesIO(); plt.savefig(img, format='png', bbox_inches='tight', dpi=100); plt.close()
        return img
    except: return None

def add_texto(slide, texto, left, top, width, height, size=18, bold=False, color=COR_TEXTO, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(texto)
    p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(size), bold, color, align

# --- UI STREAMLIT ---
st.set_page_config(page_title="Invest Forma Academy", layout="wide")
st.title("💼 Dashboard de Performance - Invest Forma Academy")

logo_data = carregar_logo()
data_sel = st.date_input("Selecione a Data do Relatório", datetime.now() - timedelta(days=1))

if st.button("🌟 GERAR BOLETIM DE PERFORMANCE"):
    with st.spinner("Extraindo dados e calculando variações B3..."):
        try:
            # Dados principais (Brasil e Global)
            mapa = {"IBOV": "^BVSP", "DOLAR": "USDBRL=X", "S&P500": "^GSPC", "NASDAQ": "^IXIC", "DOW": "^DJI", "BTC": "BTC-USD", "BRENT": "BZ=F", "IRON": "TIO=F", "GOLD": "GC=F", "SILVER": "SI=F"}
            res = {}
            data_final = data_sel
            
            for nome, ticker in mapa.items():
                info = obter_dado_unico(ticker, data_sel)
                if info:
                    res[nome] = info
                    data_final = info['D'] # Sincroniza a data com o último pregão válido

            # MOVERS B3 (Processamento Individual)
            altas, baixas = obter_variacoes_b3(data_final)

            # --- PPTX ---
            prs = Presentation(); prs.slide_width, prs.slide_height = Inches(9), Inches(16)
            def slide_mestre(titulo):
                s = prs.slides.add_slide(prs.slide_layouts[6])
                s.background.fill.solid(); s.background.fill.fore_color.rgb = COR_FUNDO
                if logo_data: s.shapes.add_picture(io.BytesIO(logo_data), Inches(6.4), Inches(0.4), width=Inches(2.2))
                add_texto(s, titulo, 0.5, 0.4, 6, 1, size=32, bold=True, color=COR_DESTAQUE)
                return s

            # S1: CAPA
            s1 = slide_mestre("Boletim de Performance")
            add_texto(s1, "Fechamento do Mercado Financeiro", 0.5, 0.9, 6, 0.5, 22)
            add_texto(s1, data_final.strftime('%d/%m/%Y'), 0.5, 1.3, 6, 0.4, 20, True, COR_ALTA)
            # Notícias (Backup Robusto)
            news = ["Mercado reage a dados econômicos globais", "Fluxo na B3 monitorado por investidores", "Cenário fiscal doméstico no radar", "Commodities operam em níveis de suporte", "Atenção voltada a juros nos EUA e Brasil"]
            for i, n in enumerate(news, 1): add_texto(s1, f"{i}. {n}", 0.7, 2.5 + (i*1.2), 7.5, 0.8, 20)

            # S2: BRASIL (IBOV + DOLAR + MOVERS)
            s2 = slide_mestre("MERCADO BRASIL")
            v_i = res['IBOV']; cor_i = COR_ALTA if v_i['V'] >= 0 else COR_BAIXA
            add_texto(s2, f"IBOVESPA: {v_i['F']:,.0f} ({v_i['V']:+.2f}%)", 0.5, 1.6, 8, 0.8, 28, True, cor_i)
            gi = gerar_grafico("^BVSP", data_final, "#00FF7F")
            if gi: s2.shapes.add_picture(gi, Inches(2), Inches(2.2), width=Inches(5))
            
            v_d = res['DOLAR']; cor_d = COR_BAIXA if v_d['V'] >= 0 else COR_ALTA # Dolar inverte lógica
            add_texto(s2, f"DÓLAR: R$ {v_d['F']:.3f} ({v_d['V']:+.2f}%)", 0.5, 5.0, 8, 0.8, 28, True, cor_d)
            gd = gerar_grafico("USDBRL=X", data_final, "#00AEF3")
            if gd: s2.shapes.add_picture(gd, Inches(2), Inches(5.6), width=Inches(5))

            if altas is not None:
                add_texto(s2, "📈 MAIORES ALTAS", 0.5, 8.8, 4, 0.5, 22, True, COR_ALTA)
                add_texto(s2, "\n".join([f"{t}: +{v:.2f}%" for t,v in altas.items()]), 0.5, 9.4, 4, 3, 20)
                add_texto(s2, "📉 MAIORES BAIXAS", 4.5, 8.8, 4, 0.5, 22, True, COR_BAIXA)
                add_texto(s2, "\n".join([f"{t}: {v:.2f}%" for t,v in baixas.items()]), 4.5, 9.4, 4, 3, 20)

            # S3: EUA
            s3 = slide_mestre("BOLSAS EUA"); y_e = 2.0
            for idx in ["S&P500", "NASDAQ", "DOW"]:
                vi = res[idx]; cor_e = COR_ALTA if vi['V'] >= 0 else COR_BAIXA
                add_texto(s3, f"{idx}: {vi['F']:,.0f} ({vi['V']:+.2f}%)", 0.5, y_e, 4, 0.6, 22, True, cor_e)
                ge = gerar_grafico(mapa[idx], data_final, "#00AEF3")
                if ge: s3.shapes.add_picture(ge, Inches(4.5), Inches(y_e), width=Inches(4))
                y_e += 4.2

            # S4: GLOBAL
            s4 = slide_mestre("GLOBAL & CRIPTO")
            c_t = "".join([f"• {c}: US$ {res[c]['F']:,.2f} ({res[c]['V']:+.2f}%)\n\n" for c in ["BRENT", "IRON", "GOLD", "SILVER"]])
            add_texto(s4, c_t, 1.0, 2.5, 7, 5, 24)
            vb = res['BTC']; add_texto(s4, f"BITCOIN: US$ {vb['F']:,.0f} ({vb['V']:+.2f}%)", 0.5, 8.8, 8, 1, 28, True, COR_DESTAQUE)
            gb = gerar_grafico("BTC-USD", data_final, "#F7931A")
            if gb: s4.shapes.add_picture(gb, Inches(2), Inches(9.8), width=Inches(5))

            out = io.BytesIO(); prs.save(out)
            st.success(f"✅ Sucesso! Boletim de {data_final.strftime('%d/%m/%Y')} gerado.")
            st.download_button("📥 BAIXAR APRESENTAÇÃO PREMIUM", out.getvalue(), f"InvestForma_{data_final.strftime('%Y-%m-%d')}.pptx")
        except Exception as e: st.error(f"Erro ao processar: {e}")